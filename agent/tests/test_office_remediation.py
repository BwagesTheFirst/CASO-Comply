"""Integration tests for Office document remediation engines."""
import os
import tempfile
from pathlib import Path

import docx
import openpyxl
from pptx import Presentation


def _create_test_docx(path: str):
    """Create a minimal Word doc with accessibility issues."""
    doc = docx.Document()
    para = doc.add_paragraph("Company Report")
    run = para.runs[0]
    run.bold = True
    run.font.size = docx.shared.Pt(24)
    doc.add_paragraph("This is the body text of the report.")
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "100"
    doc.save(path)


def _create_test_xlsx(path: str):
    """Create a minimal Excel workbook with accessibility issues."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Name"
    ws["B1"] = "Score"
    ws["A2"] = "Alice"
    ws["B2"] = 95
    ws.merge_cells("A4:B4")
    ws["A4"] = "Summary"
    wb.save(path)


def _create_test_pptx(path: str):
    """Create a minimal PowerPoint with accessibility issues."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = "Important content without a slide title"
    prs.save(path)


def test_docx_analysis():
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = f.name
    try:
        _create_test_docx(path)
        from remediation_docx import analyze_docx
        result = analyze_docx(path)
        assert "score" in result
        assert result["score"]["score"] < 100
        assert not result["has_title"]
    finally:
        os.unlink(path)


def test_xlsx_analysis():
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        path = f.name
    try:
        _create_test_xlsx(path)
        from remediation_xlsx import analyze_xlsx
        result = analyze_xlsx(path)
        assert "score" in result
        assert result["merged_cells"] > 0
        assert len(result["sheets"]["generic_names"]) > 0
    finally:
        os.unlink(path)


def test_pptx_analysis():
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    try:
        _create_test_pptx(path)
        from remediation_pptx import analyze_pptx
        result = analyze_pptx(path)
        assert "score" in result
        assert len(result["slides"]["without_title"]) > 0
    finally:
        os.unlink(path)


def test_docx_remediation():
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        path = f.name
    output = path.replace(".docx", "_remediated.docx")
    try:
        _create_test_docx(path)
        import asyncio
        from remediation_docx import remediate_docx_async
        result = asyncio.run(remediate_docx_async(path, output, verify=False))
        assert result["after"]["score"]["score"] >= result["before"]["score"]["score"]
        assert Path(output).exists()
    finally:
        for p in (path, output):
            if os.path.exists(p):
                os.unlink(p)


def test_xlsx_remediation():
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        path = f.name
    output = path.replace(".xlsx", "_remediated.xlsx")
    try:
        _create_test_xlsx(path)
        import asyncio
        from remediation_xlsx import remediate_xlsx_async
        result = asyncio.run(remediate_xlsx_async(path, output, verify=False))
        assert result["after"]["score"]["score"] >= result["before"]["score"]["score"]
        assert result["after"]["merged_cells"] == 0
    finally:
        for p in (path, output):
            if os.path.exists(p):
                os.unlink(p)


def test_pptx_remediation():
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    output = path.replace(".pptx", "_remediated.pptx")
    try:
        _create_test_pptx(path)
        import asyncio
        from remediation_pptx import remediate_pptx_async
        result = asyncio.run(remediate_pptx_async(path, output, verify=False))
        # Verify remediated file was produced with valid structure
        assert Path(output).exists()
        assert "before" in result and "after" in result
        # Presentation title should be set after remediation
        assert result["after"]["has_title"]
        # Slides without title should be fewer (generic titles added)
        assert len(result["after"]["slides"]["without_title"]) <= len(
            result["before"]["slides"]["without_title"]
        )
    finally:
        for p in (path, output):
            if os.path.exists(p):
                os.unlink(p)


def test_scanner_finds_all_formats():
    with tempfile.TemporaryDirectory() as tmpdir:
        _create_test_docx(os.path.join(tmpdir, "test.docx"))
        _create_test_xlsx(os.path.join(tmpdir, "test.xlsx"))
        _create_test_pptx(os.path.join(tmpdir, "test.pptx"))
        Path(os.path.join(tmpdir, "test.pdf")).write_bytes(b"%PDF-1.4 dummy")

        from agent.scanner import scan_folder
        results = scan_folder(tmpdir)
        formats = {r["format"] for r in results}
        assert formats == {"pdf", "docx", "xlsx", "pptx"}
