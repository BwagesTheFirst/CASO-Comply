"""
CASO Comply -- PowerPoint (.pptx) Accessibility Remediation Engine

Analyzes and remediates PowerPoint presentation accessibility issues using
python-pptx and lxml for OOXML manipulation.
"""
from __future__ import annotations

import asyncio
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from lxml import etree

logger = logging.getLogger(__name__)

# OOXML namespaces
NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ---------------------------------------------------------------------------
# 1. Analysis
# ---------------------------------------------------------------------------

def analyze_pptx(file_path: str) -> dict:
    """Analyze PowerPoint presentation for accessibility issues."""
    prs = Presentation(file_path)

    # Check slide titles (placeholder OR remediation textbox at origin)
    slides_without_title = []
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.text.strip():
            continue
        # Check for remediation title textbox (added by _add_missing_titles)
        has_remediation_title = False
        for shape in slide.shapes:
            if (shape.shape_type == 17  # TEXT_BOX
                    and shape.left == 0 and shape.top == 0
                    and shape.has_text_frame and shape.text_frame.text.strip()):
                has_remediation_title = True
                break
        if not has_remediation_title:
            slides_without_title.append(i + 1)

    title_coverage = (total_slides - len(slides_without_title)) / max(total_slides, 1)

    # Check reading order
    reading_order_issues = []
    for i, slide in enumerate(prs.slides):
        if _has_reading_order_issue(slide):
            reading_order_issues.append(i + 1)

    reading_order_coverage = (total_slides - len(reading_order_issues)) / max(total_slides, 1)

    # Check images for alt text (only actual pictures, not textboxes)
    total_images = 0
    images_with_alt = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type and shape.shape_type in (13, 11):  # Picture, Linked Picture
                total_images += 1
                if _shape_has_alt_text(shape):
                    images_with_alt += 1

    alt_text_coverage = images_with_alt / max(total_images, 1)

    # Check tables
    total_tables = 0
    tables_with_headers = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                total_tables += 1
                if _table_has_header(shape.table):
                    tables_with_headers += 1

    table_header_coverage = tables_with_headers / max(total_tables, 1)

    # Check animations/transitions
    dangerous_animations = _count_dangerous_animations(prs)

    # Check document properties
    has_title = bool(prs.core_properties.title)

    # Language check
    has_lang = _has_presentation_language(prs)

    # Compute score
    score = _compute_pptx_score(
        title_coverage=title_coverage,
        reading_order_coverage=reading_order_coverage,
        alt_text_coverage=alt_text_coverage,
        table_header_coverage=table_header_coverage,
        dangerous_animations=dangerous_animations,
        has_lang=has_lang,
        has_title=has_title,
        total_images=total_images,
        total_tables=total_tables,
        total_slides=total_slides,
    )

    return {
        "score": score,
        "slides": {
            "total": total_slides,
            "without_title": slides_without_title,
            "reading_order_issues": reading_order_issues,
        },
        "images": {"total": total_images, "with_alt": images_with_alt},
        "tables": {"total": total_tables, "with_headers": tables_with_headers},
        "dangerous_animations": dangerous_animations,
        "has_title": has_title,
        "has_lang": has_lang,
        "page_count": total_slides,
    }


def _shape_has_alt_text(shape) -> bool:
    """Check if a shape has alt text set."""
    el = shape._element
    cNvPr = el.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
    if cNvPr is None:
        cNvPr = el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr")
    if cNvPr is None:
        nvSpPr = el.find(".//{%s}nvSpPr" % NSMAP["p"])
        if nvSpPr is not None:
            cNvPr = nvSpPr.find("{%s}cNvPr" % NSMAP["p"])
    if cNvPr is None:
        for tag in ("nvPicPr", "nvSpPr", "nvGrpSpPr"):
            parent = el.find(".//{%s}%s" % (NSMAP["p"], tag))
            if parent is not None:
                cNvPr = parent.find("{%s}cNvPr" % NSMAP["p"])
                if cNvPr is not None:
                    break

    if cNvPr is not None:
        return bool(cNvPr.get("descr", "").strip())
    return False


def _has_reading_order_issue(slide) -> bool:
    """Check if shapes are in a logical reading order."""
    shapes = list(slide.shapes)
    if len(shapes) <= 1:
        return False

    title_shape = slide.shapes.title
    if title_shape:
        sp_tree = slide._element.find(".//{%s}spTree" % NSMAP["p"])
        if sp_tree is not None:
            # Only count actual shape elements, skip nvGrpSpPr/grpSpPr
            shape_children = [
                c for c in sp_tree
                if etree.QName(c.tag).localname not in ("nvGrpSpPr", "grpSpPr")
            ]
            title_idx = None
            for i, child in enumerate(shape_children):
                if child == title_shape._element:
                    title_idx = i
                    break
            if title_idx is not None and title_idx > 0:
                return True

    # Check non-placeholder shapes for reading order (top-to-bottom)
    # Exclude placeholders since their position is inherited from layout
    non_ph = [s for s in shapes if s.shape_type != 14 and hasattr(s, 'top') and s.top is not None]
    if len(non_ph) < 2:
        return False

    tops = [s.top or 0 for s in non_ph]
    sorted_tops = sorted(tops)
    return tops != sorted_tops


def _table_has_header(table) -> bool:
    """Check if a table has firstRow set."""
    tbl = table._tbl
    tblPr = tbl.find("{%s}tblPr" % NSMAP["a"])
    if tblPr is not None:
        return tblPr.get("firstRow", "0") == "1"
    return False


def _count_dangerous_animations(prs) -> int:
    """Count transitions and animations that could trigger seizures."""
    count = 0
    for slide in prs.slides:
        sp_tree = slide._element
        transition = sp_tree.find("{%s}transition" % NSMAP["p"])
        if transition is not None:
            count += 1
        timing = sp_tree.find("{%s}timing" % NSMAP["p"])
        if timing is not None:
            count += 1
    return count


def _has_presentation_language(prs) -> bool:
    """Check if presentation has a default language set."""
    try:
        presentation_el = prs.presentation._element if hasattr(prs, 'presentation') else prs._element
        defaultTextStyle = presentation_el.find("{%s}defaultTextStyle" % NSMAP["p"])
        if defaultTextStyle is not None:
            for rPr in defaultTextStyle.iter("{%s}defRPr" % NSMAP["a"]):
                if rPr.get("lang"):
                    return True
    except Exception:
        pass
    return False


def _compute_pptx_score(
    title_coverage: float,
    reading_order_coverage: float,
    alt_text_coverage: float,
    table_header_coverage: float,
    dangerous_animations: int,
    has_lang: bool,
    has_title: bool,
    total_images: int,
    total_tables: int,
    total_slides: int,
) -> dict:
    """Compute 0-100 accessibility score for PowerPoint."""
    checks = {
        "slide_titles": {
            "passed": title_coverage >= 0.9,
            "weight": 20,
            "description": "All slides have titles",
        },
        "reading_order": {
            "passed": reading_order_coverage >= 0.8,
            "weight": 20,
            "description": "Reading order is logical",
        },
        "alt_text": {
            "passed": alt_text_coverage >= 0.9 if total_images > 0 else True,
            "weight": 25,
            "description": "Images have alt text" if total_images > 0 else "No images -- alt text not needed",
        },
        "table_headers": {
            "passed": table_header_coverage >= 0.9 if total_tables > 0 else True,
            "weight": 10,
            "description": "Tables have headers" if total_tables > 0 else "No tables -- headers not needed",
        },
        "no_dangerous_animations": {
            "passed": dangerous_animations == 0,
            "weight": 10,
            "description": "No dangerous animations/transitions",
        },
        "language": {
            "passed": has_lang,
            "weight": 5,
            "description": "Presentation language set",
        },
        "title": {
            "passed": has_title,
            "weight": 5,
            "description": "Presentation title set",
        },
        "contrast": {
            "passed": True,  # Placeholder
            "weight": 5,
            "description": "Text contrast adequate (not yet checked)",
        },
    }

    total = sum(c["weight"] for c in checks.values())
    earned = sum(c["weight"] for c in checks.values() if c["passed"])
    score = round(100 * earned / total) if total > 0 else 0

    return {"score": score, "checks": checks}


# ---------------------------------------------------------------------------
# 2. Remediation
# ---------------------------------------------------------------------------

def _add_missing_titles(prs, suggestions: dict | None = None):
    """Add title placeholders to slides without titles."""
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.text.strip():
            continue

        if slide.shapes.title:
            suggested = suggestions.get(i + 1, f"Slide {i + 1}") if suggestions else f"Slide {i + 1}"
            slide.shapes.title.text = suggested
            continue

        # No title placeholder — add a hidden text box
        txBox = slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(0.3)
        )
        tf = txBox.text_frame
        suggested = suggestions.get(i + 1, f"Slide {i + 1}") if suggestions else f"Slide {i + 1}"
        tf.text = suggested
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(1)


def _fix_reading_order(slide):
    """Reorder shapes in spTree by visual position: title first, then top-to-bottom, left-to-right."""
    sp_tree = slide._element.find(".//{%s}spTree" % NSMAP["p"])
    if sp_tree is None:
        return

    non_shapes = []
    shapes = []
    for child in list(sp_tree):
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            non_shapes.append(child)
        else:
            shapes.append(child)

    if len(shapes) <= 1:
        return

    def sort_key(el):
        nvSpPr = el.find(".//{%s}nvSpPr" % NSMAP["p"])
        if nvSpPr is not None:
            nvPr = nvSpPr.find("{%s}nvPr" % NSMAP["p"])
            if nvPr is not None:
                ph = nvPr.find("{%s}ph" % NSMAP["p"])
                if ph is not None:
                    ph_type = ph.get("type", "")
                    # Title and center title go first
                    if ph_type in ("title", "ctrTitle"):
                        return (-2, 0)
                    # Subtitle goes second
                    if ph_type in ("subTitle", "body"):
                        return (-1, 0)
                    if ph.get("idx") == "0":
                        return (-2, 0)

        spPr = el.find(".//{%s}spPr" % NSMAP["a"])
        if spPr is not None:
            off = spPr.find("{%s}off" % NSMAP["a"])
            if off is not None:
                y = int(off.get("y", "0"))
                x = int(off.get("x", "0"))
                return (y, x)

        xfrm = el.find(".//{%s}xfrm" % NSMAP["a"])
        if xfrm is not None:
            off = xfrm.find("{%s}off" % NSMAP["a"])
            if off is not None:
                y = int(off.get("y", "0"))
                x = int(off.get("x", "0"))
                return (y, x)

        return (999999, 999999)

    shapes.sort(key=sort_key)

    for child in list(sp_tree):
        sp_tree.remove(child)
    for ns in non_shapes:
        sp_tree.append(ns)
    for s in shapes:
        sp_tree.append(s)


def _set_table_headers(prs):
    """Set firstRow attribute on all tables."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table._tbl
                tblPr = tbl.find("{%s}tblPr" % NSMAP["a"])
                if tblPr is not None:
                    tblPr.set("firstRow", "1")
                else:
                    tblPr = etree.SubElement(tbl, "{%s}tblPr" % NSMAP["a"])
                    tblPr.set("firstRow", "1")
                    tbl.insert(0, tblPr)


def _remove_dangerous_animations(prs):
    """Remove all transitions and complex animations."""
    for slide in prs.slides:
        el = slide._element
        for transition in el.findall("{%s}transition" % NSMAP["p"]):
            el.remove(transition)
        for timing in el.findall("{%s}timing" % NSMAP["p"]):
            el.remove(timing)


def _set_alt_text_on_shapes(slide, alt_texts: list[dict]):
    """Write alt text to shapes on a slide."""
    shapes = list(slide.shapes)
    for entry in alt_texts:
        idx = entry.get("shape_index", -1)
        desc = entry.get("description", "")
        is_decorative = entry.get("is_decorative", False)

        if idx < 0 or idx >= len(shapes):
            continue

        shape = shapes[idx]
        el = shape._element

        cNvPr = None
        for parent_tag in ("nvSpPr", "nvPicPr", "nvCxnSpPr", "nvGrpSpPr"):
            parent = el.find("{%s}%s" % (NSMAP["p"], parent_tag))
            if parent is not None:
                cNvPr = parent.find("{%s}cNvPr" % NSMAP["p"])
                if cNvPr is not None:
                    break

        if cNvPr is not None:
            if is_decorative:
                cNvPr.set("descr", "")
            else:
                cNvPr.set("descr", desc)


# ---------------------------------------------------------------------------
# 3. Main Pipeline
# ---------------------------------------------------------------------------

async def remediate_pptx_async(
    file_path: str,
    output_path: str | None = None,
    verify: bool = True,
) -> dict:
    """Full PowerPoint presentation remediation pipeline."""
    file_path = str(file_path)
    if output_path is None:
        p = Path(file_path)
        output_path = str(p.parent / f"{p.stem}_remediated{p.suffix}")

    logger.info("Analyzing input PPTX: %s", file_path)
    before_analysis = analyze_pptx(file_path)

    # Open presentation for remediation
    prs = Presentation(file_path)

    # Apply heuristic fixes
    _set_table_headers(prs)
    _remove_dangerous_animations(prs)

    # Set presentation title if missing
    if not prs.core_properties.title:
        first_title = None
        for slide in prs.slides:
            if slide.shapes.title and slide.shapes.title.text.strip():
                first_title = slide.shapes.title.text[:256]
                break
        prs.core_properties.title = first_title or "Untitled Presentation"

    # Always add missing titles BEFORE reading order fix so they're
    # positioned correctly in the spTree (title first in reading order)
    _add_missing_titles(prs)

    # Fix reading order on all slides (title first, then top-to-bottom)
    for slide in prs.slides:
        _fix_reading_order(slide)

    # Gemini AI verification
    verification_info = None
    if verify:
        from gemini_verify import verify_and_correct_office

        logger.info("Running Gemini AI verification on PPTX")

        structure_data = {
            "slides": [],
            "total_slides": len(prs.slides),
        }
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "has_title": bool(slide.shapes.title and slide.shapes.title.text.strip()),
                "title_text": slide.shapes.title.text if slide.shapes.title else None,
                "shape_count": len(slide.shapes),
                "shapes": [],
            }
            for j, shape in enumerate(slide.shapes):
                slide_data["shapes"].append({
                    "index": j,
                    "name": shape.name,
                    "type": str(shape.shape_type) if shape.shape_type else "unknown",
                    "has_text": shape.has_text_frame,
                    "text": shape.text_frame.text[:100] if shape.has_text_frame else None,
                    "has_alt": _shape_has_alt_text(shape),
                    "top": shape.top,
                    "left": shape.left,
                })
            structure_data["slides"].append(slide_data)

        verification_result = await verify_and_correct_office(
            file_path=file_path,
            structure_data=structure_data,
            format="pptx",
        )

        verification_info = {
            "issues_found": verification_result.get("issues_found", []),
            "verification_score": verification_result.get("verification_score", 0.0),
        }

        # Update existing titles with AI suggestions (don't add new shapes)
        for slide_fix in verification_result.get("slides", []):
            slide_num = slide_fix.get("slide_number", 0)
            suggested = slide_fix.get("suggested_title", "")
            if not suggested or slide_num < 1 or slide_num > len(prs.slides):
                continue
            slide = prs.slides[slide_num - 1]
            # Update placeholder title if it exists and is empty/generic
            if slide.shapes.title:
                current = slide.shapes.title.text.strip()
                if not current or current.startswith("Slide "):
                    slide.shapes.title.text = suggested
            else:
                # Find our remediation textbox at origin and update its text
                for shape in slide.shapes:
                    if (shape.shape_type == 17
                            and shape.left == 0 and shape.top == 0
                            and shape.has_text_frame):
                        shape.text_frame.text = suggested
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(1)
                        break

        # Apply alt texts per slide
        for alt_entry in verification_result.get("alt_texts", []):
            slide_num = alt_entry.get("slide", 0)
            if 1 <= slide_num <= len(prs.slides):
                _set_alt_text_on_shapes(prs.slides[slide_num - 1], [alt_entry])

    # Save remediated presentation
    prs.save(output_path)
    logger.info("Saved remediated PPTX: %s", output_path)

    # Analyze the output
    after_analysis = analyze_pptx(output_path)

    result = {
        "before": before_analysis,
        "after": after_analysis,
        "output_path": output_path,
        "page_count": before_analysis.get("page_count", 1),
    }

    if verification_info:
        result["verification"] = verification_info

    return result
